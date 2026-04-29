"""
文件上传与向量化处理脚本
独立运行，用于将 PDF / Word (.docx) 文件解析、分块并写入向量库。

特性:
  - PDF: 优先使用 PyMuPDF (fitz) 提取文本，支持用 marker 转 Markdown 保留 LaTeX 公式
  - Word: 使用 python-docx 提取文本 + OMML 数学公式转 LaTeX
  - 分块: 使用 LaTeX 公式感知的分块器，避免从中间切断 $...$ 或 $$...$$

用法:
    # 上传单个文件
    python upload_and_index.py /path/to/file.pdf

    # 上传多个文件
    python upload_and_index.py file1.pdf file2.docx file3.pdf

    # 上传整个目录（递归扫描 .pdf 和 .docx）
    python upload_and_index.py /path/to/folder/

    # 指定分块大小
    python upload_and_index.py --chunk-size 512 --chunk-overlap 100 file.pdf

    # 使用 marker 模型做 PDF→Markdown（保留 LaTeX 公式效果最好，需要 GPU）
    python upload_and_index.py --use-marker file.pdf

    # 仅预览解析结果，不写入向量库
    python upload_and_index.py --dry-run file.pdf
"""

import argparse
import os
import re
import sys
import time
from pathlib import Path

# ─────────────────────────────────────────────
# 依赖检查与安装
# ─────────────────────────────────────────────
def _ensure_deps():
    """确保文件解析依赖已安装"""
    missing = []
    try:
        import pypdf  # noqa: F401
    except ImportError:
        missing.append("pypdf")

    try:
        import fitz  # PyMuPDF  # noqa: F401
    except ImportError:
        missing.append("PyMuPDF")

    try:
        import docx2txt  # noqa: F401
    except ImportError:
        missing.append("docx2txt")

    try:
        import docx  # python-docx  # noqa: F401
    except ImportError:
        missing.append("python-docx")

    try:
        import pptx  # python-pptx  # noqa: F401
    except ImportError:
        missing.append("python-pptx")

    try:
        import openpyxl  # noqa: F401
    except ImportError:
        missing.append("openpyxl")

    try:
        import bs4  # BeautifulSoup4  # noqa: F401
    except ImportError:
        missing.append("beautifulsoup4")

    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        missing.append("Pillow")

    try:
        from openai import OpenAI  # noqa: F401 — VLM 调用需要
    except ImportError:
        missing.append("openai")

    if missing:
        print(f"📦 正在安装缺失依赖: {', '.join(missing)}")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", *missing])
        print("✅ 依赖安装完成")


_ensure_deps()

# ─────────────────────────────────────────────
# 文件加载器
# ─────────────────────────────────────────────
from langchain_core.documents import Document


# ─────────────────────────────────────────────
# 方案1: PDF — PyMuPDF + marker
# ─────────────────────────────────────────────
def load_pdf_with_fitz(file_path: str) -> list[Document]:
    """使用 PyMuPDF (fitz) 加载 PDF，文本提取质量优于 pypdf"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    docs = []
    for i, page in enumerate(doc):
        # 优先提取带排版信息的文本
        text = page.get_text("text")
        if text and text.strip():
            docs.append(Document(
                page_content=text.strip(),
                metadata={
                    "source": file_path,
                    "page": i + 1,
                    "total_pages": len(doc),
                    "file_type": "pdf",
                    "extraction_method": "pymupdf",
                }
            ))
    doc.close()
    return docs


def load_pdf_with_marker(file_path: str) -> list[Document]:
    """
    使用 marker 将 PDF 转为 Markdown（保留 LaTeX 公式）。
    需要安装: pip install marker-pdf
    需要 GPU，转换较慢但公式保留效果最好。
    """
    try:
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict
        from marker.output import text_from_rendered
    except ImportError:
        print("⚠️ marker-pdf 未安装，回退到 PyMuPDF。安装: pip install marker-pdf")
        return load_pdf_with_fitz(file_path)

    print("   ⏳ 使用 marker 转 PDF→Markdown（保留 LaTeX 公式，需要 GPU，较慢）...")
    model_dict = create_model_dict()
    converter = PdfConverter(artifact_dict=model_dict)
    rendered = converter(file_path)
    text, _, _ = text_from_rendered(rendered)

    if not text or not text.strip():
        return []

    return [Document(
        page_content=text.strip(),
        metadata={
            "source": file_path,
            "file_type": "pdf",
            "extraction_method": "marker",
            "has_latex": "$$" in text or "\\[" in text,
        }
    )]


def load_pdf(file_path: str, use_marker: bool = False) -> list[Document]:
    """加载 PDF 文件，优先使用 PyMuPDF，可选 marker"""
    if use_marker:
        try:
            return load_pdf_with_marker(file_path)
        except Exception as e:
            print(f"   ⚠️ marker 转换失败: {e}，回退到 PyMuPDF")

    return load_pdf_with_fitz(file_path)


# ─────────────────────────────────────────────
# 方案3: Word — python-docx + OMML→LaTeX
# ─────────────────────────────────────────────
def _omml_to_latex(omml_xml: str) -> str:
    """
    将 Office Math Markup Language (OMML) 转换为 LaTeX。
    
    OMML 是 Word 文档中数学公式的 XML 表示格式。
    此函数覆盖常见的 OMML 元素，将其转为 LaTeX 语法。
    对于无法识别的元素，保留原始 XML 标记。
    """
    from lxml import etree

    # 注册 OMML 命名空间
    nsmap = {
        'm': 'http://schemas.openxmlformats.org/officeDocument/2006/math',
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
    }

    try:
        root = etree.fromstring(omml_xml)
    except etree.XMLSyntaxError:
        return omml_xml

    def _convert_node(node):
        """递归转换 OMML 节点"""
        tag = etree.QName(node.tag).localname if isinstance(node.tag, str) else ''
        parts = []

        if tag == 'f':  # 分数 <m:f>
            num = ''
            den = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'num':
                    num = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'den':
                    den = ''.join(_convert_node(c) for c in child)
            return f'\\frac{{{num}}}{{{den}}}'

        elif tag == 'rad':  # 根号 <m:rad>
            deg = ''
            e = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'deg':
                    deg = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'e':
                    e = ''.join(_convert_node(c) for c in child)
            if deg:
                return f'\\sqrt[{deg}]{{{e}}}'
            return f'\\sqrt{{{e}}}'

        elif tag == 'sSup':  # 上标 <m:sSup>
            base = ''
            sup = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    base = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'sup':
                    sup = ''.join(_convert_node(c) for c in child)
            return f'{base}^{{{sup}}}'

        elif tag == 'sSub':  # 下标 <m:sSub>
            base = ''
            sub = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    base = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'sub':
                    sub = ''.join(_convert_node(c) for c in child)
            return f'{base}_{{{sub}}}'

        elif tag == 'sSubSup':  # 上下标 <m:sSubSup>
            base = ''
            sub = ''
            sup = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    base = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'sub':
                    sub = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'sup':
                    sup = ''.join(_convert_node(c) for c in child)
            return f'{base}_{{{sub}}}^{{{sup}}}'

        elif tag == 'nary':  # n-元运算 (求和、积分等) <m:nary>
            lower = ''
            upper = ''
            body = ''
            chr_val = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'sub':
                    lower = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'sup':
                    upper = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'e':
                    body = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'naryPr':
                    for prop in child:
                        prop_tag = etree.QName(prop.tag).localname if isinstance(prop.tag, str) else ''
                        if prop_tag == 'chr':
                            chr_val = prop.get(etree.QName('http://schemas.openxmlformats.org/officeDocument/2006/math', 'val'), '')

            # 根据 chr 决定 LaTeX 命令
            nary_map = {
                '∑': '\\sum', '∫': '\\int', '∏': '\\prod',
                '∬': '\\iint', '∮': '\\oint',
            }
            op = nary_map.get(chr_val, '\\sum')
            result = f'{op}'
            if lower:
                result += f'_{{{lower}}}'
            if upper:
                result += f'^{{{upper}}}'
            result += f' {body}'
            return result

        elif tag == 'd':  # 定界符 (括号等) <m:d>
            inner = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    inner = ''.join(_convert_node(c) for c in child)
            return f'\\left({inner}\\right)'

        elif tag == 'func':  # 函数名 <m:func>
            fname = ''
            body = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'fName':
                    fname = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'e':
                    body = ''.join(_convert_node(c) for c in child)
            return f'{fname}\\left({body}\\right)'

        elif tag == 'r':  # 文本运行 <m:r>
            text = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 't':
                    text += (child.text or '')
            # 处理特殊字符
            text = text.replace('∈', '\\in ').replace('≤', '\\leq ').replace('≥', '\\geq ')
            text = text.replace('≠', '\\neq ').replace('×', '\\times ').replace('÷', '\\div ')
            text = text.replace('±', '\\pm ').replace('→', '\\rightarrow ')
            text = text.replace('∞', '\\infty ').replace('∂', '\\partial ')
            text = text.replace('α', '\\alpha ').replace('β', '\\beta ').replace('γ', '\\gamma ')
            text = text.replace('δ', '\\delta ').replace('θ', '\\theta ').replace('λ', '\\lambda ')
            text = text.replace('μ', '\\mu ').replace('π', '\\pi ').replace('σ', '\\sigma ')
            text = text.replace('ω', '\\omega ').replace('φ', '\\phi ').replace('ψ', '\\psi ')
            return text

        elif tag == 'bar':  # 上划线/下划线 <m:bar>
            body = ''
            position = 'top'
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    body = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'barPr':
                    for prop in child:
                        prop_tag = etree.QName(prop.tag).localname if isinstance(prop.tag, str) else ''
                        if prop_tag == 'pos':
                            pos_val = prop.get(etree.QName('http://schemas.openxmlformats.org/officeDocument/2006/math', 'val'), 'top')
                            position = pos_val
            if position == 'bot':
                return f'\\underline{{{body}}}'
            return f'\\overline{{{body}}}'

        elif tag == 'acc':  # 重音符号 (hat, tilde, vec 等) <m:acc>
            body = ''
            chr_val = ''
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    body = ''.join(_convert_node(c) for c in child)
                elif child_tag == 'accPr':
                    for prop in child:
                        prop_tag = etree.QName(prop.tag).localname if isinstance(prop.tag, str) else ''
                        if prop_tag == 'chr':
                            chr_val = prop.get(etree.QName('http://schemas.openxmlformats.org/officeDocument/2006/math', 'val'), '')

            acc_map = {
                '̂': '\\hat', '̃': '\\tilde', '⃗': '\\vec',
                '̄': '\\bar', '̇': '\\dot', '̈': '\\ddot',
                '⏞': '\\overbrace', '⏟': '\\underbrace',
            }
            cmd = acc_map.get(chr_val, '\\hat')
            return f'{cmd}{{{body}}}'

        elif tag == 'eqArr':  # 方程组 <m:eqArr>
            rows = []
            for child in node:
                child_tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
                if child_tag == 'e':
                    rows.append(''.join(_convert_node(c) for c in child))
            return ' \\\\ '.join(rows)

        elif tag == 'm':  # 数学区域 (最外层)
            return ''.join(_convert_node(c) for c in node)

        elif tag == 'oMathPara':  # 数学段落
            return ''.join(_convert_node(c) for c in node)

        elif tag == 'oMath':  # 行内/块级数学
            return ''.join(_convert_node(c) for c in node)

        else:
            # 未知节点：递归处理子节点
            if len(node) > 0:
                return ''.join(_convert_node(c) for c in node)
            return node.text or ''

    latex = _convert_node(root)
    return latex.strip()


def _docx_cell_to_text(cell) -> str:
    """提取 Word 表格单元格的文本（含段落内 OMML 公式）"""
    from lxml import etree
    parts = []
    for para in cell.paragraphs:
        para_parts = []
        for child in para._element:
            tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
            if tag == 'r':
                for sub in child:
                    sub_tag = etree.QName(sub.tag).localname if isinstance(sub.tag, str) else ''
                    if sub_tag == 't':
                        para_parts.append(sub.text or '')
            elif tag in ('oMathPara', 'oMath'):
                try:
                    omml_str = etree.tostring(child, encoding='unicode')
                    latex = _omml_to_latex(omml_str)
                    para_parts.append(f' ${latex}$ ')
                except Exception:
                    para_parts.append(''.join(child.itertext()))
        text = ''.join(para_parts).strip()
        if text:
            parts.append(text)
    return ' '.join(parts)


def _docx_table_to_markdown(table) -> str:
    """将 Word 表格转为 Markdown 格式"""
    rows_data = []
    for row in table.rows:
        cells = [_docx_cell_to_text(cell).replace('|', '│').replace('\n', ' ') for cell in row.cells]
        # 处理合并单元格：python-docx 会对合并单元格重复返回内容，去重
        rows_data.append(cells)

    if not rows_data:
        return ''

    # 去重连续重复行（合并单元格导致）
    deduped_rows = [rows_data[0]]
    for row in rows_data[1:]:
        if row != deduped_rows[-1]:
            deduped_rows.append(row)

    col_count = max(len(r) for r in deduped_rows) if deduped_rows else 0
    if col_count == 0:
        return ''

    # 补齐列数
    for row in deduped_rows:
        while len(row) < col_count:
            row.append('')

    # 构建 Markdown 表格
    lines = []
    # 表头
    lines.append('| ' + ' | '.join(deduped_rows[0]) + ' |')
    # 分隔线
    lines.append('| ' + ' | '.join(['---'] * col_count) + ' |')
    # 数据行
    for row in deduped_rows[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(lines)


def load_docx_with_math(file_path: str) -> list[Document]:
    """
    使用 python-docx 加载 Word 文件，提取文本 + 表格 + OMML 数学公式转 LaTeX。
    
    相比 docx2txt，此方法能：
    1. 识别 Word 中的 OMML 数学公式（<m:oMath> 元素）
    2. 将 OMML 转为 LaTeX 格式，用 $...$ 包裹
    3. 保留公式的精确语义
    4. 提取表格并转为 Markdown 格式，保持表格结构可检索
    """
    from docx import Document as DocxDocument
    from lxml import etree

    doc = DocxDocument(file_path)
    content_blocks = []  # 按 Word 文档顺序存储内容块
    math_count = 0
    table_count = 0

    # 构建 body 子元素顺序映射，以便按原始顺序交替提取段落和表格
    body = doc.element.body
    para_idx = 0
    table_idx = 0

    for child in body:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''

        if tag == 'p':  # 段落
            if para_idx < len(doc.paragraphs):
                para = doc.paragraphs[para_idx]
                parts = []
                for p_child in para._element:
                    p_tag = etree.QName(p_child.tag).localname if isinstance(p_child.tag, str) else ''

                    if p_tag == 'r':  # 普通文本运行 <w:r>
                        for sub in p_child:
                            sub_tag = etree.QName(sub.tag).localname if isinstance(sub.tag, str) else ''
                            if sub_tag == 't':
                                parts.append(sub.text or '')

                    elif p_tag == 'oMathPara':  # 块级数学段落
                        for omath in p_child:
                            omath_tag = etree.QName(omath.tag).localname if isinstance(omath.tag, str) else ''
                            if omath_tag == 'oMath':
                                try:
                                    omml_str = etree.tostring(omath, encoding='unicode')
                                    latex = _omml_to_latex(omml_str)
                                    parts.append(f'\n$$\n{latex}\n$$\n')
                                    math_count += 1
                                except Exception:
                                    text = ''.join(omath.itertext())
                                    parts.append(f' ${text}$ ')
                                    math_count += 1

                    elif p_tag == 'oMath':  # 行内数学公式
                        try:
                            omml_str = etree.tostring(p_child, encoding='unicode')
                            latex = _omml_to_latex(omml_str)
                            parts.append(f' ${latex}$ ')
                            math_count += 1
                        except Exception:
                            text = ''.join(p_child.itertext())
                            parts.append(f' ${text}$ ')

                para_text = ''.join(parts).strip()
                if para_text:
                    content_blocks.append(para_text)
                para_idx += 1

        elif tag == 'tbl':  # 表格
            if table_idx < len(doc.tables):
                table = doc.tables[table_idx]
                md_table = _docx_table_to_markdown(table)
                if md_table:
                    content_blocks.append(md_table)
                    table_count += 1
                table_idx += 1

    if not content_blocks:
        return []

    full_text = '\n\n'.join(content_blocks)
    has_latex = math_count > 0

    if math_count:
        print(f"   → 提取了 {math_count} 个数学公式 (OMML→LaTeX)")
    if table_count:
        print(f"   → 提取了 {table_count} 个表格 (Markdown)")

    return [Document(
        page_content=full_text,
        metadata={
            "source": file_path,
            "file_type": "docx",
            "extraction_method": "python-docx+omml+table",
            "math_count": math_count,
            "has_latex": has_latex,
            "table_count": table_count,
        }
    )]


# 兼容旧版 docx2txt（不含公式处理）
def load_docx_simple(file_path: str) -> list[Document]:
    """使用 docx2txt 简单提取（不含公式处理）"""
    import docx2txt

    text = docx2txt.process(file_path)
    if not text or not text.strip():
        return []

    return [Document(
        page_content=text.strip(),
        metadata={
            "source": file_path,
            "file_type": "docx",
            "extraction_method": "docx2txt",
        }
    )]


def load_docx(file_path: str) -> list[Document]:
    """加载 Word 文件，优先使用 OMML→LaTeX 方式"""
    try:
        return load_docx_with_math(file_path)
    except Exception as e:
        print(f"   ⚠️ OMML 提取失败: {e}，回退到 docx2txt")
        return load_docx_simple(file_path)


def load_md(file_path: str) -> list[Document]:
    """加载 Markdown 文件，直接读取文本（Markdown 本身就用 $...$/$$...$$ 写公式）"""
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    if not text or not text.strip():
        return []

    # 检测是否包含 LaTeX 公式
    has_latex = bool(re.search(r'\$\$.+?\$\$', text, re.DOTALL)) or \
                bool(re.search(r'(?<!\$)\$(?!\$).+?(?<!\$)\$(?!\$)', text, re.DOTALL))

    # 按标题分节（保留标题作为上下文）
    sections = []
    current_section = []
    current_heading = ""

    for line in text.split("\n"):
        if line.startswith("#"):
            # 遇到新标题，保存上一节
            if current_section:
                section_text = "\n".join(current_section).strip()
                if section_text:
                    sections.append(Document(
                        page_content=section_text,
                        metadata={
                            "source": file_path,
                            "file_type": "md",
                            "heading": current_heading,
                            "has_latex": has_latex,
                        }
                    ))
            current_heading = line.strip("#").strip()
            current_section = [line]
        else:
            current_section.append(line)

    # 最后一节
    if current_section:
        section_text = "\n".join(current_section).strip()
        if section_text:
            sections.append(Document(
                page_content=section_text,
                metadata={
                    "source": file_path,
                    "file_type": "md",
                    "heading": current_heading,
                    "has_latex": has_latex,
                }
            ))

    if not sections:
        # 没有标题结构，作为单个文档
        sections = [Document(
            page_content=text.strip(),
            metadata={
                "source": file_path,
                "file_type": "md",
                "has_latex": has_latex,
            }
        )]

    return sections


def load_pptx(file_path: str) -> list[Document]:
    """
    使用 python-pptx 加载 PowerPoint 文件，逐页提取文本 + OMML 数学公式转 LaTeX。
    
    PPTX 中的数学公式与 Word 一样使用 OMML 格式，可复用 _omml_to_latex 转换器。
    """
    from pptx import Presentation
    from lxml import etree

    prs = Presentation(file_path)
    docs = []

    for i, slide in enumerate(prs.slides):
        parts = []
        math_count = 0

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                para_parts = []
                for child in para._p:  # _p 是底层 XML 元素
                    tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''

                    if tag == 'r':  # 普通文本 <a:r>
                        for sub in child:
                            sub_tag = etree.QName(sub.tag).localname if isinstance(sub.tag, str) else ''
                            if sub_tag == 't':
                                para_parts.append(sub.text or '')

                    elif tag == 'oMathPara':  # 块级数学
                        for omath in child:
                            omath_tag = etree.QName(omath.tag).localname if isinstance(omath.tag, str) else ''
                            if omath_tag == 'oMath':
                                try:
                                    omml_str = etree.tostring(omath, encoding='unicode')
                                    latex = _omml_to_latex(omml_str)
                                    para_parts.append(f'\n$$\n{latex}\n$$\n')
                                    math_count += 1
                                except Exception:
                                    text = ''.join(omath.itertext())
                                    para_parts.append(f' ${text}$ ')
                                    math_count += 1

                    elif tag == 'oMath':  # 行内数学
                        try:
                            omml_str = etree.tostring(child, encoding='unicode')
                            latex = _omml_to_latex(omml_str)
                            para_parts.append(f' ${latex}$ ')
                            math_count += 1
                        except Exception:
                            text = ''.join(child.itertext())
                            para_parts.append(f' ${text}$ ')

                para_text = ''.join(para_parts).strip()
                if para_text:
                    parts.append(para_text)

        slide_text = '\n'.join(parts).strip()
        if slide_text:
            docs.append(Document(
                page_content=slide_text,
                metadata={
                    "source": file_path,
                    "page": i + 1,
                    "total_pages": len(prs.slides),
                    "file_type": "pptx",
                    "math_count": math_count,
                    "has_latex": math_count > 0,
                }
            ))

    return docs


def load_txt(file_path: str) -> list[Document]:
    """加载纯文本文件"""
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    if not text or not text.strip():
        return []

    return [Document(
        page_content=text.strip(),
        metadata={
            "source": file_path,
            "file_type": "txt",
        }
    )]


# ─────────────────────────────────────────────
# Excel / CSV 加载器
# ─────────────────────────────────────────────
def load_xlsx(file_path: str) -> list[Document]:
    """加载 Excel 文件，逐行提取并格式化为文本"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    docs = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            # 将每行转为文本，None 跳过
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))

        if rows:
            text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
            docs.append(Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "file_type": "xlsx",
                    "sheet": sheet_name,
                }
            ))

    wb.close()
    return docs


def load_csv(file_path: str) -> list[Document]:
    """加载 CSV 文件，逐行提取并格式化为文本"""
    import csv

    docs = []
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = []
        header = None
        for i, row in enumerate(reader):
            if i == 0:
                header = row
            if row:
                rows.append(" | ".join(row))

        if rows:
            text = "\n".join(rows)
            docs.append(Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "file_type": "csv",
                    "header": header,
                }
            ))

    return docs


# ─────────────────────────────────────────────
# HTML 加载器
# ─────────────────────────────────────────────
def load_html(file_path: str) -> list[Document]:
    """加载 HTML 文件，提取正文文本"""
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")

    # 移除脚本和样式
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)

    if not text:
        return []

    # 尝试获取 title
    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()

    return [Document(
        page_content=text,
        metadata={
            "source": file_path,
            "file_type": "html",
            "title": title,
        }
    )]


# ─────────────────────────────────────────────
# JSON 加载器
# ─────────────────────────────────────────────
def load_json(file_path: str) -> list[Document]:
    """加载 JSON 文件，支持数组或对象格式"""
    import json

    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    docs = []

    if isinstance(data, list):
        # JSON 数组：每个元素作为一个 Document
        for i, item in enumerate(data):
            text = json.dumps(item, ensure_ascii=False, indent=2)
            docs.append(Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "file_type": "json",
                    "index": i,
                }
            ))
    elif isinstance(data, dict):
        # JSON 对象：尝试提取文本字段，否则整体序列化
        # 常见文本字段名
        text_fields = ["text", "content", "body", "description", "message", "question", "answer"]
        found_text = None
        for field in text_fields:
            if field in data and isinstance(data[field], str):
                found_text = data[field]
                break

        if found_text:
            docs.append(Document(
                page_content=found_text,
                metadata={
                    "source": file_path,
                    "file_type": "json",
                    **{k: str(v) for k, v in data.items() if k not in text_fields and isinstance(v, (str, int, float, bool))},
                }
            ))
        else:
            # 整体序列化
            text = json.dumps(data, ensure_ascii=False, indent=2)
            docs.append(Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "file_type": "json",
                }
            ))
    else:
        text = json.dumps(data, ensure_ascii=False, indent=2)
        docs.append(Document(
            page_content=text,
            metadata={"source": file_path, "file_type": "json"}
        ))

    return docs


# ─────────────────────────────────────────────
# LaTeX 加载器
# ─────────────────────────────────────────────
def load_latex(file_path: str) -> list[Document]:
    """
    加载 .tex 文件，按 \\section/\\subsection 分节。
    LaTeX 文件天然包含数学公式，直接保留。
    """
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    if not text or not text.strip():
        return []

    # 按 \section 分节
    sections = []
    current_heading = ""
    current_lines = []

    for line in text.split("\n"):
        # 匹配 \section{...}, \subsection{...}, \chapter{...}
        sec_match = re.match(r'\\(chapter|section|subsection|subsubsection)\{(.+?)\}', line)
        if sec_match:
            # 保存上一节
            if current_lines:
                section_text = "\n".join(current_lines).strip()
                if section_text:
                    sections.append(Document(
                        page_content=section_text,
                        metadata={
                            "source": file_path,
                            "file_type": "tex",
                            "heading": current_heading,
                            "has_latex": True,
                        }
                    ))
            current_heading = sec_match.group(2)
            current_lines = [line]
        else:
            current_lines.append(line)

    # 最后一节
    if current_lines:
        section_text = "\n".join(current_lines).strip()
        if section_text:
            sections.append(Document(
                page_content=section_text,
                metadata={
                    "source": file_path,
                    "file_type": "tex",
                    "heading": current_heading,
                    "has_latex": True,
                }
            ))

    if not sections:
        sections = [Document(
            page_content=text.strip(),
            metadata={"source": file_path, "file_type": "tex", "has_latex": True}
        )]

    return sections


# ─────────────────────────────────────────────
# EPUB 加载器
# ─────────────────────────────────────────────
def load_epub(file_path: str) -> list[Document]:
    """加载 EPUB 电子书，逐章提取文本"""
    try:
        from ebooklib import epub as epub_lib
    except ImportError:
        print("   ⚠️ ebooklib 未安装，尝试安装...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "ebooklib"])
        from ebooklib import epub as epub_lib

    book = epub_lib.read_epub(file_path)
    docs = []

    for i, item in enumerate(book.get_items_of_type(epub_lib.ITEM_DOCUMENT)):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)

        if text and text.strip():
            docs.append(Document(
                page_content=text.strip(),
                metadata={
                    "source": file_path,
                    "file_type": "epub",
                    "chapter": i + 1,
                    "chapter_id": item.get_id(),
                }
            ))

    return docs


# ─────────────────────────────────────────────
# SRT 字幕加载器
# ─────────────────────────────────────────────
def load_srt(file_path: str) -> list[Document]:
    """加载 SRT 字幕文件，合并为连续文本"""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    # SRT 格式：序号 → 时间轴 → 字幕文本 → 空行
    # 提取纯文本，去掉序号和时间轴
    lines = content.strip().split("\n")
    text_lines = []
    for line in lines:
        line = line.strip()
        # 跳过序号（纯数字行）和时间轴行
        if line.isdigit():
            continue
        if "-->" in line:
            continue
        if line:
            text_lines.append(line)

    text = " ".join(text_lines)
    if not text:
        return []

    return [Document(
        page_content=text,
        metadata={
            "source": file_path,
            "file_type": "srt",
        }
    )]


# ═══════════════════════════════════════════════
# 知识图谱格式加载器
# ═══════════════════════════════════════════════

def load_jsonld(file_path: str) -> list[Document]:
    """
    加载 JSON-LD 格式的知识图谱数据。
    
    解析 GraphRAG 导出的 JSON-LD 文件，将实体、关系和社区报告
    转换为 Document 对象，以便存入向量数据库进行语义检索。
    """
    import json as _json
    
    def _extract_name(id_uri: str) -> str:
        return id_uri.rstrip("/").split("/")[-1]
    
    def _extract_type_short(type_uri: str) -> str:
        return type_uri.rstrip("/").split("/")[-1]
    
    print(f"   📊 解析 JSON-LD 知识图谱数据...")
    
    with open(file_path, "r", encoding="utf-8") as f:
        data = _json.load(f)
    
    graph = data.get("@graph", [])
    docs = []
    
    entity_type_uris = {
        "http://schema.org/PERSON",
        "http://schema.org/ORGANIZATION",
        "http://schema.org/EVENT",
        "http://schema.org/GEO",
        "http://schema.org/INDUSTRY",
        "http://schema.org/PRODUCT",
        "http://schema.org/ORGANISM",
    }
    
    for node in graph:
        node_type = node.get("@type", "")
        node_id = node.get("@id", "")
        
        # 实体节点
        if node_type in entity_type_uris:
            name = node.get("schema:name") or node.get("rdfs:label") or _extract_name(node_id)
            entity_type_short = _extract_type_short(node_type)
            description = node.get("schema:description", "")
            frequency = node.get("kg:frequency", 0)
            degree = node.get("kg:degree", 0.0)
            
            content = f"{name} ({entity_type_short}): {description}".strip()
            docs.append(Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "data_type": "kg_entity",
                    "kg_name": name,
                    "kg_entity_type": entity_type_short,
                    "kg_frequency": int(frequency) if frequency else 0,
                    "kg_degree": float(degree) if degree else 0.0,
                    "file_type": "jsonld",
                }
            ))
        
        # 关系节点
        elif node_type == "rdf:Statement":
            subject = node.get("rdf:subject", {})
            predicate = node.get("rdf:predicate", {})
            obj = node.get("rdf:object", {})
            description = node.get("schema:description", "")
            weight = node.get("kg:weight", 1.0)
            
            subject_name = _extract_name(subject.get("@id", "")) if isinstance(subject, dict) else ""
            predicate_name = _extract_name(predicate.get("@id", "")) if isinstance(predicate, dict) else ""
            object_name = _extract_name(obj.get("@id", "")) if isinstance(obj, dict) else ""
            
            content = f"{subject_name} --[{predicate_name}]--> {object_name}: {description}".strip()
            docs.append(Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "data_type": "kg_relationship",
                    "kg_subject": subject_name,
                    "kg_predicate": predicate_name,
                    "kg_object": object_name,
                    "kg_weight": float(weight) if weight else 1.0,
                    "file_type": "jsonld",
                }
            ))
        
        # 社区报告节点
        elif node_type == "schema:Article":
            headline = node.get("schema:headline", "")
            summary = node.get("schema:description", "")
            full_content = node.get("schema:text", "")
            community_id = _extract_name(node_id)
            
            content = full_content or f"{headline}\n\n{summary}"
            docs.append(Document(
                page_content=content.strip(),
                metadata={
                    "source": file_path,
                    "data_type": "kg_community",
                    "kg_community_id": community_id,
                    "kg_headline": headline,
                    "file_type": "jsonld",
                }
            ))
    
    entity_count = sum(1 for d in docs if d.metadata.get("data_type") == "kg_entity")
    rel_count = sum(1 for d in docs if d.metadata.get("data_type") == "kg_relationship")
    comm_count = sum(1 for d in docs if d.metadata.get("data_type") == "kg_community")
    
    print(f"   → 提取了 {entity_count} 个实体, {rel_count} 个关系, {comm_count} 个社区报告")
    return docs


def load_ttl(file_path: str) -> list[Document]:
    """
    加载 Turtle (.ttl) 格式的 RDF 知识图谱数据。
    
    使用简单的正则解析，提取实体描述和关系三元组。
    对于复杂的 TTL 文件，建议先转换为 JSON-LD 格式。
    """
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    
    docs = []
    
    # 提取前缀定义
    prefixes = {}
    for match in re.finditer(r'@prefix\s+(\w*):\s*<([^>]+)>\s*\.', content):
        prefixes[match.group(1)] = match.group(2)
    
    # 按句号分割三元组
    # 简单的正则匹配: subject predicate object .
    triple_pattern = re.compile(
        r'(<[^>]+>|[\w]+:[\w]+)\s+(<[^>]+>|[\w]+:[\w]+)\s+(<[^>]+>|"([^"]*)"@?\w*|[\w]+:[\w]+)\s*\.',
        re.DOTALL
    )
    
    for match in triple_pattern.finditer(content):
        subject, predicate, obj_full, obj_literal = match.groups()
        obj = obj_literal if obj_literal else obj_full
        
        # 清理 URI
        def clean_uri(uri: str) -> str:
            if uri.startswith("<") and uri.endswith(">"):
                uri = uri[1:-1]
            return uri.rstrip("/").split("/")[-1]
        
        s_name = clean_uri(subject)
        p_name = clean_uri(predicate)
        o_name = obj.strip('"') if obj_literal else clean_uri(obj)
        
        doc_content = f"{s_name} {p_name} {o_name}"
        docs.append(Document(
            page_content=doc_content,
            metadata={
                "source": file_path,
                "data_type": "kg_triple",
                "kg_subject": s_name,
                "kg_predicate": p_name,
                "kg_object": o_name,
                "file_type": "ttl",
            }
        ))
    
    print(f"   → 提取了 {len(docs)} 个三元组")
    return docs


def load_nt(file_path: str) -> list[Document]:
    """
    加载 N-Triples (.nt) 格式的 RDF 数据。
    
    N-Triples 格式简单: <subject> <predicate> <object> .
    """
    docs = []
    
    def clean_uri(uri: str) -> str:
        uri = uri.strip()
        if uri.startswith("<") and uri.endswith(">"):
            uri = uri[1:-1]
        return uri.rstrip("/").split("/")[-1]
    
    with open(file_path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            
            # N-Triples 格式: <s> <p> <o> 或 <s> <p> "literal"@lang .
            parts = re.findall(r'(<[^>]+>|"[^"]*"(?:@\w+)?)', line)
            if len(parts) < 3:
                continue
            
            subject, predicate, obj = parts[0], parts[1], parts[2]
            
            s_name = clean_uri(subject)
            p_name = clean_uri(predicate)
            
            if obj.startswith('"'):
                # 字面值
                o_name = obj.strip('"').split('"@')[0]
            else:
                o_name = clean_uri(obj)
            
            doc_content = f"{s_name} {p_name} {o_name}"
            docs.append(Document(
                page_content=doc_content,
                metadata={
                    "source": file_path,
                    "data_type": "kg_triple",
                    "kg_subject": s_name,
                    "kg_predicate": p_name,
                    "kg_object": o_name,
                    "file_type": "nt",
                }
            ))
    
    print(f"   → 提取了 {len(docs)} 个三元组")
    return docs


# ═══════════════════════════════════════════════
# 图片加载器 — OCR + VLM 双通道
# ═══════════════════════════════════════════════

def _ocr_with_paddleocr(image_path: str, lang: str = "ch") -> str:
    """使用 PaddleOCR 提取图片文字"""
    try:
        from paddleocr import PaddleOCR
    except ImportError:
        print("   ⚠️ PaddleOCR 未安装，尝试安装...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "paddleocr", "paddlepaddle"])
        from paddleocr import PaddleOCR

    ocr = PaddleOCR(use_angle_cls=True, lang=lang, show_log=False)
    result = ocr.ocr(image_path, cls=True)
    if not result or not result[0]:
        return ""
    lines = [line[1][0] for line in result[0] if line[1][0].strip()]
    return "\n".join(lines)


def _ocr_with_easyocr(image_path: str, lang: str = "ch") -> str:
    """使用 EasyOCR 提取图片文字"""
    try:
        import easyocr
    except ImportError:
        print("   ⚠️ EasyOCR 未安装，尝试安装...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "easyocr"])
        import easyocr

    lang_list = ["ch_sim", "en"] if lang == "ch" else [lang]
    reader = easyocr.Reader(lang_list)
    result = reader.readtext(image_path)
    lines = [text for (_, text, _) in result if text.strip()]
    return "\n".join(lines)


def _ocr_with_tesseract(image_path: str, lang: str = "chi_sim") -> str:
    """使用 Tesseract OCR 提取图片文字"""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        print("   ⚠️ pytesseract/Pillow 未安装，尝试安装...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "pytesseract", "Pillow"])
        import pytesseract
        from PIL import Image

    tesseract_lang = "chi_sim+eng" if lang == "ch" else lang
    img = Image.open(image_path)
    return pytesseract.image_to_string(img, lang=tesseract_lang).strip()


def _ocr_image(image_path: str) -> str:
    """
    自动选择 OCR 引擎提取图片文字。
    优先级: paddleocr > easyocr > tesseract
    """
    from config import OCR_ENGINE, OCR_LANGUAGE

    engine = OCR_ENGINE
    lang = OCR_LANGUAGE

    if engine == "paddleocr":
        return _ocr_with_paddleocr(image_path, lang)
    elif engine == "easyocr":
        return _ocr_with_easyocr(image_path, lang)
    elif engine == "tesseract":
        return _ocr_with_tesseract(image_path, lang)

    # auto 模式：按优先级尝试
    errors = []
    for name, func in [("paddleocr", _ocr_with_paddleocr), ("easyocr", _ocr_with_easyocr), ("tesseract", _ocr_with_tesseract)]:
        try:
            text = func(image_path, lang)
            if text.strip():
                return text
        except Exception as e:
            errors.append(f"{name}: {e}")
            continue

    if errors:
        print(f"   ⚠️ 所有 OCR 引擎均失败: {'; '.join(errors)}")
    return ""


def _vlm_describe_image(image_path: str) -> str:
    """
    使用 VLM (视觉语言模型) 描述图片内容。
    支持通义千问 Qwen-VL / OpenAI GPT-4o 等。
    """
    from config import VLM_BACKEND, VLM_MODEL, VLM_API_KEY, VLM_BASE_URL
    import base64

    # 读取图片并编码为 base64
    with open(image_path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode("utf-8")

    ext = Path(image_path).suffix.lower()
    mime_map = {
        ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".png": "image/png", ".gif": "image/gif",
        ".bmp": "image/bmp", ".webp": "image/webp",
        ".tiff": "image/tiff", ".tif": "image/tiff",
    }
    mime_type = mime_map.get(ext, "image/jpeg")

    prompt = """请详细描述这张图片的内容。包括：
1. 图片中的所有文字内容（逐字抄录）
2. 图片的类型（如：发票、合同、表格、照片、图表、印章、签名等）
3. 图片中的关键视觉元素（颜色、布局、图形等）
4. 如果是文档类图片，说明文档的格式和结构

请用中文回答，尽量详细和准确。"""

    if VLM_BACKEND in ("tongyi", "openai"):
        try:
            from openai import OpenAI

            client = OpenAI(api_key=VLM_API_KEY, base_url=VLM_BASE_URL)
            response = client.chat.completions.create(
                model=VLM_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{img_b64}"},
                            },
                        ],
                    }
                ],
                max_tokens=1024,
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"   ⚠️ VLM 描述失败: {e}")
            return ""

    # 本地 VLM（如 Qwen2-VL）
    elif VLM_BACKEND == "local":
        try:
            from transformers import AutoProcessor, AutoModelForVision2Seq
            import torch

            # 懒加载
            if not hasattr(_vlm_describe_image, "_local_model"):
                model_name = VLM_MODEL
                print(f"   🔧 加载本地 VLM: {model_name}...")
                _vlm_describe_image._local_processor = AutoProcessor.from_pretrained(model_name)
                _vlm_describe_image._local_model = AutoModelForVision2Seq.from_pretrained(
                    model_name, torch_dtype=torch.float16, device_map="auto"
                )

            from PIL import Image
            img = Image.open(image_path).convert("RGB")
            processor = _vlm_describe_image._local_processor
            model = _vlm_describe_image._local_model

            inputs = processor(text=prompt, images=img, return_tensors="pt").to(model.device)
            output = model.generate(**inputs, max_new_tokens=512)
            return processor.decode(output[0], skip_special_tokens=True).strip()
        except Exception as e:
            print(f"   ⚠️ 本地 VLM 描述失败: {e}")
            return ""

    return ""


def load_image(file_path: str) -> list[Document]:
    """
    加载图片文件，使用 OCR + VLM 双通道提取内容。

    - OCR 通道：提取图片上的文字
    - VLM 通道：视觉模型描述图片内容（印章、图表、照片等无文字内容）
    - 双通道合并去重 → 用 BGE-M3 编码文本 → 存入 Milvus
    """
    import shutil
    from config import IMAGE_STORAGE_DIR

    print(f"🖼️ 加载图片: {file_path}")
    print("   🔍 双通道处理: OCR + VLM...")

    # 复制图片到存储目录（检索时可返回图片路径）
    os.makedirs(IMAGE_STORAGE_DIR, exist_ok=True)
    stored_path = os.path.join(IMAGE_STORAGE_DIR, os.path.basename(file_path))
    if not os.path.exists(stored_path):
        shutil.copy2(file_path, stored_path)

    ocr_text = ""
    vlm_text = ""

    # 通道 1: OCR
    try:
        ocr_text = _ocr_image(file_path)
        if ocr_text:
            print(f"   ✅ OCR 提取文字: {len(ocr_text)} 字符")
        else:
            print(f"   ℹ️ OCR 未检测到文字（可能是无文字图片）")
    except Exception as e:
        print(f"   ⚠️ OCR 处理失败: {e}")

    # 通道 2: VLM
    try:
        vlm_text = _vlm_describe_image(file_path)
        if vlm_text:
            print(f"   ✅ VLM 描述: {len(vlm_text)} 字符")
        else:
            print(f"   ℹ️ VLM 未返回描述")
    except Exception as e:
        print(f"   ⚠️ VLM 处理失败: {e}")

    # 合并两通道内容
    content_parts = []

    if ocr_text and vlm_text:
        content_parts.append(f"【图片文字内容】\n{ocr_text}")
        content_parts.append(f"【图片视觉描述】\n{vlm_text}")
    elif ocr_text:
        content_parts.append(f"【图片文字内容】\n{ocr_text}")
    elif vlm_text:
        content_parts.append(f"【图片视觉描述】\n{vlm_text}")
    else:
        content_parts.append(f"图片文件: {os.path.basename(file_path)}")

    full_content = "\n\n".join(content_parts)

    return [Document(
        page_content=full_content,
        metadata={
            "source": file_path,
            "stored_image_path": stored_path,
            "file_type": "image",
            "data_type": "image",
            "ocr_text_length": len(ocr_text),
            "vlm_text_length": len(vlm_text),
            "has_ocr_text": bool(ocr_text),
            "has_vlm_description": bool(vlm_text),
            "extraction_method": "ocr+vlm",
        }
    )]


# 文件类型 → 加载函数 映射
LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".md": load_md,
    ".markdown": load_md,
    ".txt": load_txt,
    ".text": load_txt,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".csv": load_csv,
    ".tsv": load_csv,
    ".html": load_html,
    ".htm": load_html,
    ".json": load_json,
    ".jsonl": load_json,
    ".jsonld": load_jsonld,     # 知识图谱 JSON-LD
    ".ttl": load_ttl,           # 知识图谱 Turtle
    ".nt": load_nt,             # 知识图谱 N-Triples
    ".tex": load_latex,
    ".latex": load_latex,
    ".epub": load_epub,
    ".srt": load_srt,
    ".vtt": load_srt,  # WebVTT 字幕格式，与 SRT 类似
    # 图片格式 — OCR + VLM 双通道
    ".jpg": load_image, ".jpeg": load_image,
    ".png": load_image, ".gif": load_image,
    ".bmp": load_image, ".webp": load_image,
    ".tiff": load_image, ".tif": load_image,
}


def load_file(file_path: str, use_marker: bool = False) -> list[Document]:
    """根据扩展名自动选择加载器"""
    ext = Path(file_path).suffix.lower()
    loader = LOADERS.get(ext)
    if loader is None:
        print(f"⚠️ 不支持的文件类型: {ext}  (支持: {', '.join(LOADERS.keys())})")
        return []
    print(f"📄 加载文件: {file_path}")

    if ext == ".pdf":
        docs = loader(file_path, use_marker=use_marker)
    else:
        docs = loader(file_path)

    # 统计公式信息
    has_latex = any(d.metadata.get("has_latex", False) for d in docs)
    math_total = sum(d.metadata.get("math_count", 0) for d in docs)
    if has_latex or math_total:
        print(f"   → 包含数学公式，LaTeX 格式保留")

    print(f"   → 提取了 {len(docs)} 页/段")
    return docs


def scan_directory(dir_path: str) -> list[str]:
    """递归扫描目录下的 PDF 和 DOCX 文件"""
    files = []
    for root, _, filenames in os.walk(dir_path):
        for fname in filenames:
            if Path(fname).suffix.lower() in LOADERS:
                files.append(os.path.join(root, fname))
    return files


# ─────────────────────────────────────────────
# 方案2: LaTeX 公式感知分块器
# ─────────────────────────────────────────────
class LatexAwareTextSplitter:
    """
    LaTeX 公式感知的文本分块器。
    
    在 RecursiveCharacterTextSplitter 基础上增加：
    1. 将 $$...$$ 和 $...$ 包裹的公式识别为不可分割单元
    2. 分块时优先在公式边界处断开，避免从公式中间截断
    3. 若单个公式超过 chunk_size，整块保留（不做切割）
    """

    # 公式正则：匹配 $$...$$ 和 $...$
    DISPLAY_MATH_RE = re.compile(r'\$\$(.+?)\$\$', re.DOTALL)
    INLINE_MATH_RE = re.compile(r'(?<!\$)\$(?!\$)(.+?)(?<!\$)\$(?!\$)', re.DOTALL)

    def __init__(self, chunk_size: int = 1024, chunk_overlap: int = 200, **kwargs):
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        # LaTeX 公式感知的分隔符优先级：
        # 1. 双换行（段落边界）
        # 2. 块级公式边界 $$...$$
        # 3. 单换行
        # 4. 中文句号等标点
        # 5. 行内公式边界 $...$
        # 6. 空格
        self._splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            model_name="gpt-4",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            add_start_index=True,
            separators=[
                "\n\n",           # 段落
                "$$",             # 块级公式边界
                "\n",             # 换行
                "。", "！", "？",  # 中文标点
                "；", "，",
                "$",              # 行内公式边界
                " ",              # 空格
                "",               # 最后兜底
            ],
        )

    def split_documents(self, documents: list[Document]) -> list[Document]:
        """分块文档，保护公式不被截断"""
        result = []
        for doc in documents:
            text = doc.page_content
            # 先保护公式：用占位符替换，避免被分块器切断
            protected_text, formulas = self._protect_formulas(text)

            # 分块
            splits = self._splitter.split_text(protected_text)

            # 还原公式
            for split_text in splits:
                restored = self._restore_formulas(split_text, formulas)
                # 跳过过短的块
                if len(restored.strip()) < 10:
                    continue
                result.append(Document(
                    page_content=restored.strip(),
                    metadata={**doc.metadata},
                ))

        return result

    def _protect_formulas(self, text: str) -> tuple[str, dict[str, str]]:
        """
        用占位符替换公式，防止分块器从公式中间截断。
        
        Returns:
            (protected_text, formula_map)
            protected_text: 用占位符替换后的文本
            formula_map: 占位符 → 原始公式的映射
        """
        formulas = {}
        counter = [0]

        def replace_display(m):
            key = f"__FORMULA_D{counter[0]}__"
            formulas[key] = m.group(0)  # 保留 $$
            counter[0] += 1
            return key

        def replace_inline(m):
            key = f"__FORMULA_I{counter[0]}__"
            formulas[key] = m.group(0)  # 保留 $
            counter[0] += 1
            return key

        # 先替换块级公式 $$...$$，再替换行内公式 $...$
        protected = self.DISPLAY_MATH_RE.sub(replace_display, text)
        protected = self.INLINE_MATH_RE.sub(replace_inline, protected)

        return protected, formulas

    def _restore_formulas(self, text: str, formulas: dict[str, str]) -> str:
        """将占位符还原为原始公式"""
        for key, formula in formulas.items():
            text = text.replace(key, formula)
        return text


# ─────────────────────────────────────────────
# 主流程
# ─────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(
        description="上传 PDF/Word 文件到自适应 RAG 向量库（支持数学公式提取）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "paths",
        nargs="+",
        help="文件或目录路径（目录会递归扫描 .pdf 和 .docx）",
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=None,
        help="分块大小（默认使用 config.py 中的 CHUNK_SIZE）",
    )
    parser.add_argument(
        "--chunk-overlap",
        type=int,
        default=None,
        help="分块重叠大小（默认使用 config.py 中的 CHUNK_OVERLAP）",
    )
    parser.add_argument(
        "--collection",
        type=str,
        default=None,
        help="Milvus 集合名（默认使用 config.py 中的 COLLECTION_NAME）",
    )
    parser.add_argument(
        "--use-marker",
        action="store_true",
        help="使用 marker 模型做 PDF→Markdown（保留 LaTeX 公式效果最好，需要 GPU）",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="仅解析文件，不写入向量库（用于预览）",
    )
    args = parser.parse_args()

    # ── 1. 收集文件 ────────────────────────
    all_files = []
    for p in args.paths:
        if os.path.isdir(p):
            found = scan_directory(p)
            print(f"📂 扫描目录 {p}: 找到 {len(found)} 个文件")
            all_files.extend(found)
        elif os.path.isfile(p):
            all_files.append(p)
        else:
            print(f"⚠️ 路径不存在: {p}")

    if not all_files:
        print("❌ 没有找到可处理的文件")
        sys.exit(1)

    print(f"\n📋 待处理文件: {len(all_files)} 个")
    for f in all_files:
        print(f"   - {f}")

    # ── 2. 解析文件 ────────────────────────
    print("\n" + "=" * 50)
    print("阶段 1/3: 解析文件（含数学公式提取）")
    print("=" * 50)

    all_docs = []
    for f in all_files:
        try:
            docs = load_file(f, use_marker=args.use_marker)
            all_docs.extend(docs)
        except Exception as e:
            print(f"❌ 解析失败 {f}: {e}")
            import traceback
            traceback.print_exc()

    if not all_docs:
        print("❌ 没有提取到任何文本内容")
        sys.exit(1)

    total_chars = sum(len(d.page_content) for d in all_docs)
    math_docs = [d for d in all_docs if d.metadata.get("has_latex") or d.metadata.get("math_count", 0) > 0]
    print(f"\n✅ 解析完成: {len(all_docs)} 页/段, 共 {total_chars:,} 字符")
    if math_docs:
        total_formulas = sum(d.metadata.get("math_count", 0) for d in math_docs)
        print(f"   📐 其中 {len(math_docs)} 页/段包含数学公式, 共 {total_formulas} 个公式已转为 LaTeX")

    # ── 3. 分块（公式感知）─────────────────
    from config import CHUNK_SIZE, CHUNK_OVERLAP, COLLECTION_NAME

    chunk_size = args.chunk_size or CHUNK_SIZE
    chunk_overlap = args.chunk_overlap or CHUNK_OVERLAP
    collection_name = args.collection or COLLECTION_NAME

    print("\n" + "=" * 50)
    print("阶段 2/3: 文本分块（LaTeX 公式感知）")
    print("=" * 50)

    # 检测是否包含 LaTeX 公式，自动选择分块器
    has_any_latex = any(
        d.metadata.get("has_latex") or "$$" in d.page_content or (d.page_content.count("$") >= 2)
        for d in all_docs
    )

    if has_any_latex:
        print("   🔍 检测到 LaTeX 公式，使用公式感知分块器")
        splitter = LatexAwareTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
    else:
        print("   📝 未检测到公式，使用标准分块器")
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )

    doc_splits = splitter.split_documents(all_docs)
    print(f"✅ 分块完成: {len(doc_splits)} 个文档块 (chunk_size={chunk_size}, overlap={chunk_overlap})")

    if args.dry_run:
        print("\n🏁 --dry-run 模式，跳过向量库写入。前 3 个文档块预览:")
        for i, doc in enumerate(doc_splits[:3]):
            preview = doc.page_content[:300].replace("\n", "\\n")
            print(f"\n   [{i+1}] source={doc.metadata.get('source','?')} page={doc.metadata.get('page','?')}")
            print(f"       {preview}...")
        return

    # ── 4. 写入向量库 ──────────────────────
    print("\n" + "=" * 50)
    print("阶段 3/3: 写入向量库")
    print("=" * 50)

    from document_processor import DocumentProcessor

    print("⏳ 初始化 DocumentProcessor（加载嵌入模型，可能需要 10-30 秒）...")
    t0 = time.time()
    doc_processor = DocumentProcessor()
    print(f"✅ DocumentProcessor 初始化完成 ({time.time()-t0:.1f}s)")

    if not doc_processor.vectorstore:
        doc_processor.initialize_vectorstore()

    # 检查哪些文件已经索引过（去重）
    existing_files = doc_processor.check_existing_urls(all_files)
    new_files = [f for f in all_files if f not in existing_files]
    skipped = len(all_files) - len(new_files)
    if skipped:
        print(f"⏭️ 跳过 {skipped} 个已存在的文件（重复）")
    if not new_files:
        print("✅ 所有文件已存在，无需重复索引")
        return

    # 只保留新文件的分块结果
    new_sources = set(new_files)
    new_doc_splits = [d for d in doc_splits if d.metadata.get("source") in new_sources]

    if not new_doc_splits:
        print("✅ 所有文件已存在，无需重复索引")
        return

    print(f"⏳ 正在向量化并写入 {len(new_doc_splits)} 个文档块到集合 '{collection_name}'...")
    t1 = time.time()
    doc_processor.add_documents_to_vectorstore(new_doc_splits)
    elapsed = time.time() - t1
    print(f"✅ 写入完成! ({elapsed:.1f}s, {len(new_doc_splits)/elapsed:.1f} docs/s)")

    # ── 5. 验证 ────────────────────────────
    print("\n" + "=" * 50)
    print("验证")
    print("=" * 50)
    try:
        retriever = doc_processor.vectorstore.as_retriever(search_kwargs={"k": 2})
        test_query = new_doc_splits[0].page_content[:50]
        results = retriever.invoke(test_query)
        print(f"🔍 用测试查询检索到 {len(results)} 个结果")
        if results:
            print(f"   前 50 字: {results[0].page_content[:50]}...")
    except Exception as e:
        print(f"⚠️ 验证检索失败: {e}")

    print(f"\n🏁 全部完成! 共处理 {len(new_files)} 个新文件, 生成 {len(new_doc_splits)} 个文档块")


if __name__ == "__main__":
    main()
